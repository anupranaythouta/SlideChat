import logging
from collections import defaultdict
from pathlib import Path

log = logging.getLogger(__name__)
_converter = None


def _get_converter():
    global _converter
    if _converter is None:
        from docling.document_converter import DocumentConverter
        _converter = DocumentConverter()
    return _converter


def _iter_picture_shapes(shapes):
    """Recursively yield picture-bearing shapes, including those nested in groups."""
    for shape in shapes:
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            yield from _iter_picture_shapes(shape.shapes)
        elif hasattr(shape, "image"):
            yield shape


def _ocr_pptx_images_by_index(file_path: str, empty_indices: set[int]) -> dict[int, str]:
    """OCR embedded images on the given 1-based PPTX slide indices."""
    extracted: dict[int, str] = {}
    try:
        from pptx import Presentation
        from rapidocr import RapidOCR
        from PIL import Image
        import io
    except ImportError as e:
        log.warning("OCR dependencies unavailable: %s", e)
        return extracted

    try:
        ocr = RapidOCR()
        prs = Presentation(file_path)
    except Exception as e:
        log.warning("OCR/PPTX init failed: %s", e)
        return extracted

    for idx, slide in enumerate(prs.slides, start=1):
        if idx not in empty_indices:
            continue
        texts: list[str] = []
        for shape in _iter_picture_shapes(slide.shapes):
            try:
                img = Image.open(io.BytesIO(shape.image.blob))
                if img.mode not in ("RGB", "L"):
                    img = img.convert("RGB")
            except Exception as e:
                log.debug("slide %d: PIL could not open image: %s", idx, e)
                continue

            try:
                # RapidOCR 3.x returns RapidOCROutput; older versions return (result, elapse)
                raw = ocr(img)
                result = raw.txts if hasattr(raw, "txts") else raw[0]
            except Exception as e:
                log.debug("slide %d: OCR call failed: %s", idx, e)
                continue

            if not result:
                continue

            # RapidOCR 3.x: txts is a flat tuple of strings
            # older API: result is list of [bbox, text, score]
            for entry in result:
                if isinstance(entry, str):
                    if entry.strip():
                        texts.append(entry.strip())
                elif len(entry) >= 2 and entry[1] and entry[1].strip():
                    texts.append(entry[1].strip())

        if texts:
            extracted[idx] = "\n".join(texts)
            log.info("slide %d: OCR recovered %d text fragments", idx, len(texts))

    return extracted


def parse_file(file_path: str) -> list[dict]:
    result = _get_converter().convert(file_path)
    doc = result.document

    page_content: dict[int, list[str]] = defaultdict(list)

    for item in doc.texts:
        if item.prov and item.text.strip():
            page_content[item.prov[0].page_no].append(item.text.strip())

    for item in doc.tables:
        if item.prov:
            try:
                md = item.export_to_markdown()
                if md.strip():
                    page_content[item.prov[0].page_no].append(md.strip())
            except Exception:
                pass

    is_pptx = Path(file_path).suffix.lower() in (".pptx", ".ppt")

    if is_pptx:
        # For PPTX, ignore Docling's page ordering and trust python-pptx's slide order.
        # OCR every slide that ended up with no text from Docling, by python-pptx index.
        try:
            from pptx import Presentation
            n_slides = len(Presentation(file_path).slides)
        except Exception:
            n_slides = len(doc.pages) if doc.pages else len(page_content)

        # Map Docling pages -> 1..n by sorted page_no, same as before
        all_pages = sorted(doc.pages.keys()) if doc.pages else sorted(page_content.keys())
        slides = []
        for i in range(1, n_slides + 1):
            pno = all_pages[i - 1] if i - 1 < len(all_pages) else None
            text = "\n".join(page_content.get(pno, [])) if pno is not None else ""
            slides.append({"slide_number": i, "text": text})

        empty = {s["slide_number"] for s in slides if not s["text"].strip()}
        if empty:
            ocr_hits = _ocr_pptx_images_by_index(file_path, empty)
            for s in slides:
                if s["slide_number"] in ocr_hits:
                    s["text"] = ocr_hits[s["slide_number"]]
    else:
        all_pages = sorted(doc.pages.keys()) if doc.pages else sorted(page_content.keys())
        slides = [
            {"slide_number": i, "text": "\n".join(page_content.get(pno, []))}
            for i, pno in enumerate(all_pages, start=1)
        ]

    return slides
